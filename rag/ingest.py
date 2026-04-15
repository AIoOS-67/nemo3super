"""Ingest docs from ../docs/ into Chroma. Supports .txt, .md, .pdf."""
import sys
from pathlib import Path
import chromadb
import csv
from pypdf import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from embedder import embed

ROOT = Path(__file__).parent.parent
DOCS_DIR = ROOT / "docs"
DB_DIR = ROOT / "chroma_db"
CHUNK_SIZE = 400
CHUNK_OVERLAP = 60


def read_file(path: Path) -> str:
    suf = path.suffix.lower()
    if suf == ".pdf":
        return "\n".join(p.extract_text() or "" for p in PdfReader(path).pages)
    if suf == ".docx":
        doc = Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    if suf == ".xlsx":
        wb = load_workbook(str(path), data_only=True, read_only=True)
        lines = []
        for ws in wb.worksheets:
            lines.append(f"# Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    lines.append("\t".join(cells))
        return "\n".join(lines)
    if suf == ".csv":
        lines = []
        with open(path, newline="", encoding="utf-8", errors="ignore") as f:
            for row in csv.reader(f):
                if any(c.strip() for c in row):
                    lines.append("\t".join(row))
        return "\n".join(lines)
    if suf == ".pptx":
        prs = Presentation(str(path))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"# Slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        txt = "".join(r.text for r in para.runs).strip()
                        if txt:
                            lines.append(txt)
        return "\n".join(lines)
    return path.read_text(encoding="utf-8", errors="ignore")


def chunk(text: str) -> list[str]:
    chunks = []
    i = 0
    while i < len(text):
        chunks.append(text[i : i + CHUNK_SIZE])
        i += CHUNK_SIZE - CHUNK_OVERLAP
    return [c.strip() for c in chunks if c.strip()]


def main():
    DOCS_DIR.mkdir(exist_ok=True)
    exts = {".txt", ".md", ".pdf", ".py", ".json", ".html",
            ".docx", ".xlsx", ".csv", ".pptx"}
    skip_dirs = {".venv", "venv", "node_modules", "__pycache__", ".git",
                 "site-packages", "dist-info", ".pytest_cache"}
    files = [
        p for p in DOCS_DIR.rglob("*")
        if p.suffix.lower() in exts
        and not any(part in skip_dirs for part in p.parts)
    ]
    print(f"Found {len(files)} file(s)\n")
    if not files:
        print(f"[!] No documents. Drop files into {DOCS_DIR} and re-run.")
        return

    client = chromadb.PersistentClient(path=str(DB_DIR))
    col = client.get_or_create_collection("knowledge_base")

    for f in files:
        print(f"📄 {f.relative_to(ROOT)}")
        text = read_file(f)
        chunks = chunk(text)
        if not chunks:
            continue
        # batch embed (NVIDIA API supports batches, cap at 50 for safety)
        for b in range(0, len(chunks), 50):
            batch = chunks[b : b + 50]
            vecs = embed(batch, input_type="passage")
            ids = [f"{f.name}::{b+i}" for i in range(len(batch))]
            col.upsert(ids=ids, documents=batch, embeddings=vecs,
                       metadatas=[{"source": str(f.name)}] * len(batch))
        print(f"   ✅ {len(chunks)} chunks")
    print(f"\nTotal entries in collection: {col.count()}")


if __name__ == "__main__":
    main()
